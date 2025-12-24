import os
import shutil
import numpy as np
import urllib.request
from pathlib import Path
from typing import List, Generator, Tuple

from constants import CHAT_MODEL_PATH, CHAT_MODEL_URL, CHAT_MODEL_FILENAME, EMBED_MODEL_NAME, MODEL_DIR

# --- AI IMPORTS ---
try:
    # Fast, Native Embeddings (The Speed Upgrade)
    from sentence_transformers import SentenceTransformer
    
    # Clustering
    from sklearn.cluster import AgglomerativeClustering
    
    # Local LLM for Naming (The GGUF Model)
    from llama_cpp import Llama # type: ignore
    
    # Text Extraction
    import pypdf
    from docx import Document
    try:
        from pptx import Presentation
    except ImportError:
        Presentation = None
    try:
        import openpyxl
    except ImportError:
        openpyxl = None
    
    AI_AVAILABLE = True
except ImportError as e:
    print(f"AI Import Error: {e}")
    AI_AVAILABLE = False


# Global Instances (Singleton Pattern)
_embed_instance = None
_chat_instance = None

def get_embed_model():
    """Loads the fast SentenceTransformer model."""
    global _embed_instance
    if not AI_AVAILABLE: raise ImportError("AI modules not loaded.")
    
    if _embed_instance is None:
        # This automatically downloads the ~90MB model to your local cache
        # It is highly optimized for CPU.
        _embed_instance = SentenceTransformer(EMBED_MODEL_NAME) #type: ignore
    return _embed_instance

def get_chat_model():
    """Loads the GGUF model for naming folders."""
    global _chat_instance
    if not AI_AVAILABLE: raise ImportError("AI modules not loaded.")
    
    if _chat_instance is None:
        if not CHAT_MODEL_PATH.exists():
            raise FileNotFoundError(f"Chat model missing: {CHAT_MODEL_PATH}")
        
        # Load GGUF Model
        _chat_instance = Llama( #type: ignore
            model_path=str(CHAT_MODEL_PATH),
            n_ctx=2048,
            verbose=False,
            n_threads=4
        )
    return _chat_instance

def scan_files(path: str) -> List[Path]:
    files = []
    try:
        for entry in os.scandir(path):
            if entry.is_file():
                files.append(Path(entry.path))
    except PermissionError:
        pass
    return files

def extract_text(file_path: Path) -> str:
    """Robust text extraction for multiple file formats."""
    text = ""
    suffix = file_path.suffix.lower()
    
    try:
        if suffix in ['.txt', '.md', '.py', '.js', '.c', '.cpp', '.h', '.java', '.json', '.xml', '.yml', '.sql', '.sh']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
        elif suffix == '.csv':
             with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                lines = [f.readline() for _ in range(30)] # Read first 30 lines
                text = "\n".join(lines)
        elif suffix == '.pdf':
            with open(file_path, 'rb') as f:
                reader = pypdf.PdfReader(f) #type: ignore
                # Read max 3 pages to save time
                for page in reader.pages[:3]:
                    extracted = page.extract_text()
                    if extracted: text += extracted + "\n"
        elif suffix == '.docx':
            doc = Document(file_path) #type: ignore
            for i, para in enumerate(doc.paragraphs):
                if i > 50: break
                text += para.text + "\n"
        elif suffix == '.pptx' and Presentation:
            prs = Presentation(file_path) #type: ignore
            for i, slide in enumerate(prs.slides):
                if i > 5: break 
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n" #type: ignore
        elif suffix == '.xlsx' and openpyxl:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            ws = wb.active
            for i, row in enumerate(ws.iter_rows(values_only=True)): #type: ignore
                if i > 20: break
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                text += row_text + "\n"
            wb.close()
        else:
            # Fallback for generic text files
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read(1000)
                    # Simple heuristic: is it mostly alphanumeric?
                    if sum(c.isalnum() for c in content) > len(content) * 0.3:
                        text = content
            except: pass
        
        return text[:4000] # Limit context size
    except Exception as e:
        print(f"Error reading {file_path.name}: {e}")
        return ""

def generate_embedding(content: str) -> List[float]:
    """Generates high-quality vectors using Sentence-Transformers."""
    if not AI_AVAILABLE: raise ImportError("AI modules not loaded")
    
    try:
        model = get_embed_model()
        clean_content = content[:1000].replace('\n', ' ').strip()
        
        if not clean_content: 
            return [0.0] * 384 # MiniLM dimension is 384
            
        # .encode returns a numpy array, convert to list
        vector = model.encode(clean_content).tolist()
        return vector
            
    except Exception as e:
        print(f"Embedding Gen Error: {e}")
        return [0.0] * 384

def cluster_embeddings(embeddings_list: List[List[float]]):
    """
    Clusters using Agglomerative Clustering with Dynamic Threshold.
    This automatically decides 'How many groups' based on content.
    """
    if not AI_AVAILABLE: raise ImportError("AI modules not loaded")
    
    n_samples = len(embeddings_list)
    if n_samples == 0: return []
    if n_samples == 1: return [0]

    try:
        np_embeddings = np.array(embeddings_list, dtype=np.float32)
        
        # Normalize vectors (Important for Cosine Similarity / Euclidean distance)
        norms = np.linalg.norm(np_embeddings, axis=1, keepdims=True)
        norms[norms == 0] = 1
        normalized_embeddings = np_embeddings / norms
        
        # Use Distance Threshold (1.5 is a good balance for MiniLM)
        # Lower = More small specific folders
        # Higher = Fewer giant generic folders
        clustering = AgglomerativeClustering( #type: ignore
            n_clusters=None, # Auto-detect number of clusters
            distance_threshold=1.5, 
            metric='euclidean', 
            linkage='ward'
        )
        
        return clustering.fit_predict(normalized_embeddings)
        
    except Exception as e:
        print(f"Clustering error: {e}")
        # Fallback: put everyone in group 0
        return [0] * n_samples

def get_smart_folder_name(files_in_cluster: List[Path], file_texts: List[str]) -> str:
    """Categorize files using the Local LLM (Llama/Qwen)."""
    if not AI_AVAILABLE: return "Group"
    try:
        llm = get_chat_model()

        # Create a prompt with previews of 5 files
        file_previews = []
        for i, f in enumerate(files_in_cluster[:5]): 
            content_preview = file_texts[i][:150].replace('\n', ' ') 
            file_previews.append(f"- {f.name}: {content_preview}...")

        input_data = "\n".join(file_previews)
        
        system_prompt = (
            "You are a file organizer. "
            "Task: Generate a short, concise folder name (max 3 words) for these files. "
            "Rules: No punctuation. Use Underscores. PascalCase. No sentences. No explanation. No generic names like 'Files'."
            "If unsure, output 'Documents'.\n\n"
        )
        user_prompt = f"Files:\n{input_data}\n\nFolder Name:"

        response = llm.create_chat_completion(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.1, # Low temp for deterministic naming
            max_tokens=15
        )
        
        content = response['choices'][0]['message']['content'].strip() # type: ignore
        
        # Cleanup
        content = content.replace("Folder Name:", "").replace('"', '').strip()
        clean_name = "".join([c for c in content if c.isalnum() or c in ('_', '-')])
        
        if len(clean_name) > 25: clean_name = clean_name[:25]
        return clean_name if clean_name else "Misc_Docs"
        
    except Exception as e:
        print(f"Naming error: {e}")
        return "Group"

def check_local_model_ready() -> Tuple[bool, List[str]]:
    """Checks if the CHAT model is present. Embeddings download auto."""
    if not AI_AVAILABLE: return False, ["AI Dependencies missing"]
    
    missing = []
    # We only check for the Chat Model (GGUF)
    # The Embedding model is handled by SentenceTransformer cache
    if not CHAT_MODEL_PATH.exists():
        missing.append(CHAT_MODEL_FILENAME)
        
    return len(missing) == 0, missing

def download_local_model(model_filename: str) -> Generator[str, None, None]:
    """Downloads the GGUF model with progress updates."""
    MODEL_DIR.mkdir(parents=True, exist_ok=True)
    
    url = CHAT_MODEL_URL
    dest = CHAT_MODEL_PATH
    
    yield f"⬇️ Downloading Chat Model..."
    yield f"   Source: {url}"

    try:
        # We use urlopen with chunked reading to yield progress
        with urllib.request.urlopen(url) as response:
            total_size = int(response.info().get('Content-Length', -1))
            
            with open(dest, 'wb') as f:
                downloaded = 0
                block_size = 8192 * 8  # 64KB buffer
                last_percent = -1
                
                while True:
                    buffer = response.read(block_size)
                    if not buffer:
                        break
                    
                    f.write(buffer)
                    downloaded += len(buffer)
                    
                    # Calculate progress
                    if total_size > 0:
                        percent = int(downloaded * 100 / total_size)
                        # Yield updates every 1% to prevent UI flooding
                        if percent > last_percent:
                            mb_down = downloaded / (1024 * 1024)
                            mb_total = total_size / (1024 * 1024)
                            # The 'Downloading:' prefix is KEY for app.py to overwrite the line
                            yield f"Downloading: {percent}% ({mb_down:.1f}MB / {mb_total:.1f}MB)"
                            last_percent = percent
                    else:
                        # Fallback if Content-Length header is missing
                        if downloaded % (1024 * 1024 * 5) == 0: # Every 5MB
                            mb_down = downloaded / (1024 * 1024)
                            yield f"Downloading: {mb_down:.1f}MB"

        if dest.exists() and dest.stat().st_size > 1000:
             yield f"✅ Download Complete: {model_filename}"
        else:
             yield f"❌ Download failed (Empty file)"

    except Exception as e:
        if dest.exists(): dest.unlink()
        yield f"❌ Connection Error: {str(e)}"